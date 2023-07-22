import unittest
from pptx import Presentation
from pptx_maker import readDataFile, addTitleSlide, addTextSlide, addListSlide, addImgSlide, addPlotSlide, makePresentation, main, isValidImage
import json
from tempfile import NamedTemporaryFile
import tempfile
from PIL import Image
import os

class TestPptxMaker(unittest.TestCase):

    def test_AddTitleSlide(self):
        presentation = Presentation()
        addTitleSlide(presentation, "Title", "Subtitle")
        slides = presentation.slides
        self.assertEqual(len(slides),1)
        slide = slides[0]  
        title = slide.shapes.title.text
        subtitle = slide.placeholders[1].text
        self.assertEqual(title, "Title")
        self.assertEqual(subtitle, "Subtitle")

    def test_AddTextSlide(self):
        presentation = Presentation()
        addTextSlide(presentation, "Title", "Long text content")
        slides = presentation.slides
        self.assertEqual(len(slides),1)
        slide = slides[0]  
        title = slide.shapes.title.text
        content = slide.placeholders[1].text
        self.assertEqual(title, "Title")
        self.assertEqual(content, "Long text content")

    def test_IsValidImage(self):
        # We test the isValidImage function for existent and nonexistent files and existent but not image files
        
        # First we test whether the function works correctly for existing image files
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp:
            temp_name = temp.name
            image = Image.new('RGB', (1, 1))
            image.save(temp_name)

        result_existing = isValidImage(temp_name)
        self.assertTrue(result_existing)

        os.remove(temp_name)

        # We deleted the tempfile above, so we can use its filepath to test the case of a nonexistent file
        result_nonexistent = isValidImage(temp_name)
        self.assertFalse(result_nonexistent)

        # Finally we test the case when the file exists but is not an image file
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as temp:
            temp_name = temp.name
            temp.write(b"Test text")
            
        result_text_file = isValidImage(temp_name)
        self.assertFalse(result_text_file)

        # Cleanup
        os.remove(temp_name)


    def test_ReadDataFile(self):

        # First, we test if the function can read correctly formatted data properly
        valid_data = [
            (1.0, 2.0),
            (3.0, 4.0),
            (5.6, 6.7)
        ]
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            for line in valid_data:
                string_line = [str(value) for value in line]
                temp_file.write((';'.join(string_line) + '\n').encode())
        temp_file_name = temp_file.name  
        self.assertTrue(readDataFile(temp_file_name))
        os.remove(temp_file_name)

        # Next, we test if the function throws the appropriate exception if an entry in the data is not a float
        string_data = [
            (1.0, "text"),
            (3.0, 4.0),
            (5.6, 6.7)
        ]
        with NamedTemporaryFile(delete=False) as temp_file:
            for line in string_data:
                string_line = [str(value) for value in line]
                temp_file.write((';'.join(string_line) + '\n').encode())
        temp_file_name = temp_file.name  
        with self.assertRaises(ValueError):
            readDataFile(temp_file.name)
        os.remove(temp_file_name)


        # Finally, we test if the function throws the appropriate exception if a line in the data is not two numbers separated by a semicolon
        long_data = [
            (1.0, 2.0, 6.0),
            (3.0, 4.0),
            (5.6, 6.7)
        ]
        with NamedTemporaryFile(delete=False) as temp_file:
            for line in long_data:
                string_line = [str(value) for value in line]
                temp_file.write((';'.join(string_line) + '\n').encode())
        temp_file_name = temp_file.name  
        with self.assertRaises(ValueError):
            readDataFile(temp_file.name)
        os.remove(temp_file_name)

    def test_AddListSlide(self): # Result can be evaluated visually
        presentation = Presentation()
        json_string = '[{ "level" : 1, "text" : "The Level 1 Text"},{ "level" : 2, "text" : "The Level 2 Text"},{ "level" : 2, "text" : "The Level 2 Text"},{ "level" : 1, "text" : "The Level 1 Text"}]'
        json_data = json.loads(json_string)
        addListSlide(presentation, "Title", json_data)
        slides = presentation.slides
        self.assertEqual(len(slides),1)
        slide = slides[0]  
        title = slide.shapes.title.text
        self.assertEqual(title, "Title")


    def test_AddImgSlide(self):  # The result of the tested function can be evaluated visually
        presentation = Presentation()
        addImgSlide(presentation, "Title", "picture.png")
        slides = presentation.slides
        self.assertEqual(len(slides),1)
        slide = slides[0]  
        title = slide.shapes.title.text
        self.assertEqual(title, "Title")
        presentation.save("img_test_pres.pptx")

    def test_AddPlotSlide(self): # The result of this function is also to be evaluated visually
        presentation = Presentation()
        addPlotSlide(presentation, "Title", "sample.dat", "xlabel", "ylabel")
        slides = presentation.slides
        self.assertEqual(len(slides),1)
        slide = slides[0]  
        title = slide.shapes.title.text
        self.assertEqual(title, "Title")
        presentation.save("plot_test_pres.pptx")

    """
    What follows is a test suite for the makePresentation function.
    Multiple functions are employed to test the normal course of the function, as well as the numerous exceptions that can be raised by the function.
    """
    
    def test_MakePresentation(self): # A visual test for the expected functioning of the makePresentation()
        with open("sample.json") as json_file:
            presentation_data = json.load(json_file)
        presentation = makePresentation(presentation_data)
        presentation.save("total_test_pres.pptx")

    def test_MissingPresentationKey(self): # Testing the case where the JSON structure has no "presentation" key
        invalid_json = {}  
        with self.assertRaises(KeyError):
            makePresentation(invalid_json)
    
    def test_MissingRequiredKeys(self): # Missing "type" key
        invalid_json = {"presentation": [{"title": "Title", "content": "Content"}]}  
        with self.assertRaises(KeyError):
            makePresentation(invalid_json)
    
    def test_InvalidSlideType(self):
        invalid_json = {"presentation": [{"type": "invalid", "title": "Title", "content": "Content"}]}  # Invalid "type" attribute
        with self.assertRaises(ValueError):
            makePresentation(invalid_json)
    
    def test_MissingConfigurationKey(self):
        invalid_json = {"presentation": [{"type": "plot", "title": "Title", "content": "Content"}]}  # Missing "configuration" key in plot slide
        with self.assertRaises(KeyError):
            makePresentation(invalid_json)


if __name__ == "__main__":
    unittest.main()

def deleteAllDemonstrations(): # A function for manually deleting the visual demonstrations made by some test functions above
    
    files_to_delete = ["img_test_pres.pptx", "plot_test_pres.pptx", "total_test_pres.pptx"]

    for filename in files_to_delete:
        try:
            os.remove(filename)
            print(f"Deleted {filename}")
        except FileNotFoundError:
            print(f"{filename} not found")
