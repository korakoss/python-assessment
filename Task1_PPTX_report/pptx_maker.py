from pptx import Presentation
from pptx.util import Inches, Pt
import json
import matplotlib.pyplot as plt
import tempfile
import os
import logging
from PIL import Image


def addTitleSlide(presentation, title_text, subtitle_text):
    """
    Adds a title slide to a given presentation, with the specified title and subtitle.

    Args:
        presentation (Presentation): an instance of the pptx Presentation class to which the title slide will be added
        title_text (str): the text to be used as the title of the slide
        subtitle_text (str): the text to be used as the subtitle of the slide

    Returns:
        None
    """
    
    slide_layout = presentation.slide_layouts[0]  
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def addTextSlide(presentation, title_text, text):
    """
    Adds a text slide to a given presentation, with the specified title and body text.

    Args:
        presentation (Presentation): an instance of the pptx Presentation class to which the text slide will be added
        title_text (str): the title text of the text slide
        text (str): the body text to be included in the slide

    Returns:
        None
    """
    
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = title_text
    body.text = text

def addListSlide(presentation, title_text, list_json):
    """
    Adds a slide with a multi-level list to a given presentation, with a specified title and list structure.

    Args:
        presentation (pptx.Presentation): An instance of the pptx.Presentation class to add the list slide to.
        title_text (str): The title text of the slide.
        list_json (list of dict): A list of dictionaries where each dictionary represents a list item on the slide. Each dictionary has keys 'level' and 'text'.
          The 'level' key should indicate the indentation level of the list item (integer values starting from 1), and 'text' key should 
          contain the text content of the list item.

    Raises:
        ValueError: If the 'level' attribute in any list item is not greater than 0 or is not an integer.

    Returns:
        None
    """
    
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    content = slide.placeholders[1]
    content.text = ""
    for list_item in list_json:
        level = int(list_item["level"])
        if not level>0:
            raise ValueError(f"Invalid 'level' attribute in JSON entry {list_item} in the list slide titled {title_text}")
        text = list_item["text"]
        paragraph = content.text_frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level - 1 

def isValidImage(filepath):
    """
    Checks if a given file path points to a valid image file or not. 

    Args:
        filepath (str): the file path of the file to check

    Returns:
        bool: True if the file is a valid image, False otherwise.
    """
    
    try:
        with Image.open(filepath) as img:
            img.verify()
        return True
    except IOError:
        return False

def addImgSlide(presentation, title_text, img_path):
    """
    Adds an image slide to a given presentation, with the specified title and image content.

    Args:
        presentation (Presentation): an instance of the pptx Presentation class to which the image slide will be added
        title_text (str): the title text of the image slide
        img_path (str): the file path of the image file to be included on the slide
        
    Raises:
        ValueError: If the img_path arg is not the file path of a valid image file.
        
    Returns:
        None
    """
    
    slide_layout = presentation.slide_layouts[5] 
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    if not isValidImage(img_path):
        raise ValueError(f"The filepath {img_path} included in your JSON file does not point to a valid image file.")
        
    slide.shapes.add_picture(img_path, Inches(1), Inches(2))

def readDataFile(filepath): 
    """
    Reads pairs of numbers from a data file, if the file is correctly formatted.

    Parameters:
        filepath (str): the file path of the data file 

    Raises:
        ValueError: If the file is not correctly formatted, ie. does not consist of lines with two floating point numbers in each line, separated by a semicolon.
        
    Returns:
        list of tuples: The list of pairs of floating point numbers read from the file.
    """
    
    data = []
    with open(filepath, 'r') as data_file:
        for line in data_file:
            line = line.strip()
            if line:
                values = line.split(';')
                if len(values) == 2:
                    try:
                        x_value = float(values[0])
                        y_value = float(values[1])
                        data.append((x_value, y_value))
                    except ValueError:
                        raise ValueError(f"Warning: Problem with line {line} in data file {filepath}.")
                else:
                    raise ValueError(f"Line {line} in data file {filepath} incorrectly formatted.")
    return data

def createPlotImage(datapoints, x_label, y_label):
    """
    Creates a matplotlib line plot that fits on a set of data points, with specified labels on the axes of the plot.
    (Based on the description of the exercise and the provided example input and output, I assumed that plot data consists of pairs of floats, and the task is to draw a line plot that fits on these points.)
    The plot is then saved to a tempfile.

    Parameters:
        datapoints (list of tuples): a list of 2-tuples of floating point numbers. The line plot will be fitted to these points.
        x_label (str): the text of the label of the x-axis of the plot
        y_label (str): the text of the label of the y-axis of the plot

    Raises:
        IOError: If the image file cannot be saved.
        
    Returns:
        str: The file name of the tempfile that contains the image of the plot.
    """
    plt.figure(figsize=(6,4))
    datapoints.sort()
    plt.plot([x for x, y in datapoints], [y for x, y in datapoints], marker='o') 
    plt.xlabel(x_label)
    plt.ylabel(y_label)
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    plt.savefig(temp_file.name)
    plt.close()
    temp_file.close()
    return temp_file.name

def addChartToPlotSlide(slide, datapoints, x_label, y_label):
    """
    Creates a matplotlib line plot and adds it as an image to a given slide. The plot is created using specified data points and axis labels.

    Parameters:
        slide (Slide): an instance of the pptx Slide class to which the plot image will be added
        datapoints (list of tuples): a list of 2-tuples of floating point numbers. The line plot will be fitted to these points.
        x_label (str): the text of the label of the x-axis of the plot
        y_label (str): the text of the label of the y-axis of the plot

    Raises:
        IOError: If the tempfile cannot be opened (or saved by the createPlotImage() called by this function)
        OSError: If the tempfile cannot be removed.
        
    Returns:
        None
    """
    
    image_path = createPlotImage(datapoints, x_label, y_label)
    slide.shapes.add_picture(image_path, Inches(1), Inches(2))
    os.remove(image_path)

def addPlotSlide(presentation, title_text, data_path, x_label, y_label):
    """
    Adds a plot slide to a presentation using provided title, plot data, and axis labels.

    Args:
        presentation (pptx.Presentation): An instance of the pptx Presentation class to which the plot slide will be added.
        title_text (str): The title text of the plot slide.
        data_path (str): The path to the file that contains the data for the plot.
        x_label (str): The label for the x-axis of the plot.
        y_label (str): The label for the y-axis of the plot.

    Raises:
        IOError: If there are issues with reading the data file or creating the plot image.
        OSError: If there are issues with removing the created plot image file.
        
    Returns:
        None
    """
    
    slide_layout = presentation.slide_layouts[5] 
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    datapoints = readDataFile(data_path)            
    addChartToPlotSlide(slide, datapoints, x_label, y_label)

def makePresentation(json_data):
    """
    Creates a pptx Presentation based on a given JSON structure.

    The JSON structure should be a dictionary with the key "presentation" that points to a list of dictionaries. Each dictionary in the list represents a slide in the presentation, and should contain at least three keys: "type", "title", and "content".
    The "type" key determines the type of the slide (it can be "title", "text", "list", "picture", or "plot"). The "title" key specifies the title of the slide. The "content" key varies depending on the type of the slide, it can be text, a list of text items, a file path of an image, or a file path of a data file.
    In the case of "plot" type slides, an additional key "configuration" is required in the dictionary of the slide. This key should point to another dictionary with the keys "x-label" and "y-label" to label the axes of the plot.

    Args:
        json_data (dict): The JSON structure that defines the presentation.

    Raises:
        KeyError: If the JSON structure is missing any required keys.
        ValueError: If anything in the JSON structure is assigned an invalid value (e.g. slide types, list levels)

    Returns:
        pptx.Presentation: The presentation that was created based on the input JSON data.
    """
    
    presentation = Presentation()
    try:
        json_root = json_data["presentation"]
    except KeyError:
        raise KeyError("Your JSON file has no top level object named 'presentation', despite it being required. Please check your JSON file.")
    
    for slide_data in json_root:
        try:
            slide_type = slide_data["type"]
            slide_title = slide_data["title"]
            slide_content = slide_data["content"]
        except KeyError as c:
            raise KeyError(f"In the JSON file, slide object {slide_data} had no key {str(c)}, despite it being required. Please check your JSON file.") 
        
        if slide_type == "title":
            addTitleSlide(presentation, slide_title, slide_content)
        elif slide_type == "text":
            addTextSlide(presentation, slide_title, slide_content)
            
        elif slide_type == "list":
            addListSlide(presentation, slide_title, slide_content)
                
        elif slide_type == "picture":
            addImgSlide(presentation, slide_title, slide_content)

        elif slide_type == "plot":
            try:
                slide_config = slide_data["configuration"]
                x_label = slide_config["x-label"]
                y_label = slide_config["y-label"]
            except KeyError as c:
                raise KeyError(f"In the JSON file, the plot slide object {slide_data} had no key {c}, despite it being required. Please check your JSON file.")
            addPlotSlide(presentation, slide_title, slide_content, x_label, y_label)

        else:
            raise ValueError(f"Incorrect slide type attribute ({slide_type}) given for the slide object {slide_data}")

    return presentation




# This is the main section of the code
# The user is first prompted for the name of a JSON file
# The JSON file is read and the data is used to create a PowerPoint presentation
# Then, the user is prompted again for the filename of the resulting pptx file. The program attempts to save the presentation under that filename
# Important events during the execution of the program are logged to 'pptx_maker.log'
# If there are errors at any stage of the process, they are logged, and the user is prompted to try again

def main():
    
    # Basic configuration of the logging system
    logging.basicConfig(filename='pptx_maker.log', filemode='a', format='%(name)s - %(levelname)s - %(message)s', level=logging.INFO)

    # Inform the user about the program
    print("This program makes presentations from JSON files.")
    print("You will be asked to provide the JSON file to be summarized as a presentation.")
    print("Make sure that the JSON file and other relevant files are in the program library.")

    while True:
        # Request the JSON file, then try to read it. If any error is encountered, the user is prompted to try again
        filename_input = input(f"\n Enter the filename of the JSON file (without the .json extension) to be turned into a PPT: ")
        json_filename = filename_input + ".json"
        logging.info(f"Attempting to read {json_filename}.")
        try:
            with open(json_filename, 'r') as file:
                presentation_data = json.load(file)
        except FileNotFoundError:
            print(f"File {json_filename} not found. Please enter a valid filename. Restarting with a new file input request.")
            logging.error(f"File {json_filename} not found.")
            continue
        
        except json.JSONDecodeError as e:
            print(f"There was an issue interpreting your JSON file. Make sure the file is valid then start the process over. Restarting with a new file input request.")
            logging.error(f"JSON decoding error with {json_filename}.")
            continue

        # Try to convert the JSON data into a presentation. If any error is encountered, the user is informed about the error, then the program starts over from requesting the JSON file
        try:
            presentation = makePresentation(presentation_data)        

        except FileNotFoundError as e:
            missing_file = e.filename
            print(f"The input file {missing_file} mentioned in your JSON source file was not found. Please check this then start the process over. Restarting with a new file input request.")
            logging.error(f"Source file {missing_file} not found.")
            continue

        except ValueError as e:
            print(f"There were issues with your JSON file or other data files. Please check the files then start the process over. Error details: {str(e)} Restarting with a new file input request.")
            logging.error(f"Value error encountered. Error message: {str(e)}")
            continue

        except KeyError as e:
            print(f"Your JSON file were missing required data. Please check the file, then start the process over. Error details: {str(e)} Restarting with a new file input request.")
            logging.error(f"Key error encountered. Error message: {str(e)}")
            continue

        except PermissionError as e:
            print(f"The program was denied permission to access file {e.filename}. Please make sure this program has the appropriate permissions, then start the process over. Restarting with a new file input request.")
            logging.error(f"Permission error encountered when trying to access {e.filename}")
            continue

        except TypeError as e:
            print(f"There were issues with some data you provided. Error details: {str(e)}. Please revise the mentioned data then start the process over. Restarting with a new file input request.")
            logging.error(f"Key error encountered. Error message: {str(e)}")
            continue

        print("Your JSON file has been successfully converted to a presentation.")

        # The user is asked to provide a filename for the resulting pptx file. The program attempts to save it under that filename. If an error is encountered, the program starts over from requesting the JSON file
        try:
            ppt_filename_input = input(f"Please enter a filename for the .pptx file to be created from your JSON file: ")
            output_filename = ppt_filename_input + ".pptx"
            presentation.save(output_filename)
            logging.info(f"JSON input file {json_filename} succesfully converted, resulting presentation saved to {output_filename}.")
            input(f"The presentation has been saved into the file {output_filename}. Enter anything to exit the program.")
            break
        
        except PermissionError:
            print(f"There was a permission error when trying to save the presentation. Please make sure this program has the appropriate permissions, then start the process over. Restarting with a new file input request.")
            logging.error(f"Permission error when trying to save the presentation to {output_filename}.")
            continue

if __name__ == "__main__":
    main()


